import os
import re
from pptx import Presentation
from flask import current_app
import logging
logger = logging.getLogger("templateR") # gets module-specific logger





def count_files():
	return len([f for f in os.listdir(current_app.config['UPLOAD_FOLDER']) if f.endswith('.pptx')]),

def get_files():
	return [f for f in os.listdir(current_app.config['UPLOAD_FOLDER']) if f.endswith('.pptx') and os.path.isfile(os.path.join(current_app.config['UPLOAD_FOLDER'], f))]

def get_routes(app):
    output = []
    for rule in app.url_map.iter_rules():
        methods = ','.join(rule.methods - {'HEAD', 'OPTIONS'})
        output.append({'endpoint': rule.endpoint, 'methods': methods,'route': str(rule)})
    return output

def slide_count(template_path):
    template_path = os.path.join(current_app.config['UPLOAD_FOLDER'], template_path)
    try: prs = Presentation(template_path) ; slide_count = len(prs.slides)
    except Exception: slide_count = None
    return slide_count

def get_metadata(template):
    path = os.path.join(current_app.config['UPLOAD_FOLDER'], template)
    metadata = {}
    prs   = Presentation(path)
    props = prs.core_properties
    return {
            'Title'            : props.title,
            'Author'           : props.author,
            'Subject'          : props.subject,
            'Created'          : props.created,
            'Modified'         : props.modified,
            'Last Modified By' : props.last_modified_by,
            'Revision'         : props.revision
        }

def placeholder_mapper(template_path):
    # Returns:
    # - slide_tag_map: {slide index: [tags]}
    # - slides_without_tags: [slide indices]
    # - all_tags: sorted list of unique tags

    prs     = Presentation(template_path)
    pattern = re.compile(r"\{\{(.*?)\}\}")

    slide_tag_map = {}
    all_tags_set = set()
    all_slide_indices = set(range(len(prs.slides)))

    for slide_index, slide in enumerate(prs.slides):
        tags_on_slide = set()
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        matches = pattern.findall(run.text)
                        tags_on_slide.update(matches)
        if tags_on_slide:
            slide_tag_map[slide_index] = sorted(tags_on_slide)
            all_tags_set.update(tags_on_slide)

    slides_with_tags = set(slide_tag_map.keys())
    slides_without_tags = sorted(all_slide_indices - slides_with_tags)
    all_tags = sorted(all_tags_set)

    return dict(sorted(slide_tag_map.items())), slides_without_tags, all_tags

def fill_placeholders(template_path, replacements):
    prs = Presentation(template_path)
    pattern = re.compile(r"\{\{(.*?)\}\}")

    def replace_tag(match):
        tag = match.group(1)
        full_tag = match.group(0)
        value = replacements.get(tag)

        if value == '!!!':
            logger.info(f"API Fill: tag retained: {tag} ")
            return full_tag  # leave tag unchanged
        elif value == "":
            logger.info(f"API Fill: tag blanked: {tag} ")
            return full_tag  # leave tag unchanged
        else:
            return value

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = pattern.sub(replace_tag, run.text)

    return prs


def filter_slides(prs, keep_indices, logger=None):
    """
    Removes slides from a presentation object that are not in keep_indices.
    
    Args:
        prs: pptx.Presentation object
        keep_indices: iterable of slide indices to keep
        logger: optional logging object
    """
    try:
        keep_set = set(map(int, keep_indices))
        total = len(prs.slides)
        for i in reversed(range(total)):
            if i not in keep_set:
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
        if logger:
            logger.info(f"Filtered slides: kept {len(prs.slides)} of {total}")
    except Exception as e:
        if logger:
            logger.error(f"Slide filtering failed: {e}")
        raise
