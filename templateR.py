'''
ROUGH PoC 7/10/25 NMT
PPT Template Filler
'''
from dotenv import load_dotenv
load_dotenv()

import os,re, json
import logging
import tempfile

from datetime import datetime
from flask import Flask, request, jsonify, session, send_file, url_for, after_this_request, render_template, flash, redirect
from flask import current_app

from werkzeug.utils import secure_filename
from functools import wraps

from pptx import Presentation

from custom_logging import setup_logging

app = Flask(__name__)
setup_logging(app)
from filler import count_files, get_files, get_routes, slide_count, get_metadata, placeholder_mapper, fill_placeholders, filter_slides


##### CONFIGURATION ############################################################
app.secret_key 						= os.getenv('SECRET_KEY')
app.config['TENANT'] 				= os.getenv('TENANT')
app.config['REMOVE_FILLED_FILES'] 	= os.getenv('REMOVE_FILLED_FILES') == 'True'
app.config['UPLOAD_FOLDER'] 		= os.path.join(app.root_path, os.getenv('UPLOAD_FOLDER'))
app.config['FILLED_FOLDER'] 		= os.getenv('FILLED_FOLDER')
app.config['ADMIN_USERNAME'] 		= os.getenv('ADMIN_USERNAME')
app.config['ADMIN_PASSWORD'] 		= os.getenv('ADMIN_PASSWORD')
app.config['API_KEY'] 				= os.getenv('API_KEY')


##### END CONFIGURATION #########################################################


os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['FILLED_FOLDER'], exist_ok=True)

app.logger.info("Starting Server")



#######################################################
#####  CONTEXT PROCESSORS                         #####
#######################################################
@app.context_processor
def inject_globals():
    return {'tenant': app.config.get('TENANT', 'unknown')}


#######################################################
#####  DECORATORS                                 #####
#######################################################
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not session.get('logged_in'):
            route = request.path  
            current_app.logger.error(f"Login Required: Access denied to {route}")
            flash("Please log in to access this page.")
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function


def require_api_key(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        key      = request.headers.get('X-API-Key') or request.args.get('api_key')
        expected = current_app.config.get('API_KEY')
        route    = request.path  
        if not key or key != expected:
            current_app.logger.error(f"API Key Required: Unauthorized {route}")
            return jsonify({'error': 'Unauthorized'}), 401
        return f(*args, **kwargs)
    return decorated



#######################################################
#####  AUTHENTICATION                             #####
#######################################################

@app.route('/ui/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')

        if username == app.config['ADMIN_USERNAME'] and password == app.config['ADMIN_PASSWORD']:
            session['logged_in'] = True
            flash("Login successful")
            app.logger.info("Successful Login")
            return redirect(url_for('list_templates_web'))
        else:
            flash("Invalid credentials")
            app.logger.warning("Failed Login: Invalid Credentials")
            return redirect(url_for('login'))

    return render_template('login.html',now=datetime.now())

@app.route('/ui/logout')
def logout():
    session.pop('logged_in', None)
    flash("You have been logged out.")
    app.logger.info("User Logged Out")
    return redirect(url_for('login'))


#######################################################
#####  OBSERVABILITY                              #####
#######################################################
## Probably a database to count fills by template

#######################################################
#####  UI PAGES                                   #####
#######################################################

@app.route('/ui', methods=['GET'])
@login_required
def ui_redirect(): return redirect(url_for('list_templates_web'))



@app.route('/ui/list-templates-web', methods=['GET'])
@login_required
def list_templates_web():
    files = get_files() ; templates = []

    for f in files:
        templates.append({
            'filename'         : f,
            'slide_count'      : slide_count(f),
            'placeholders_url' : url_for('ui_placeholders', template=f),
            'download_url'     : url_for('download_template', template=f, api_key=app.config['API_KEY']),
            'fill_url'         : url_for('fill_template_web', template=f),

        })

    status_url = url_for('status', api_key=app.config['API_KEY'], _external=True)
    return render_template('list_templates.html', templates = templates, now=datetime.now(), status_url = status_url)


@app.route('/ui/upload-web', methods=['GET', 'POST'])
@login_required
def upload_template_web():
    if request.method == 'POST':
        file = request.files.get('file')
        if file and file.filename.endswith('.pptx'):
            save_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(save_path)
            msg = f"UI Uploaded Template: {file.filename}"
            flash(msg) ; app.logger.info(msg)
            return redirect(url_for('upload_template_web'))
        else:
            flash("Please upload a valid .pptx file")
            app.logger.warning('UI Upload Failed')
            return redirect(url_for('upload_template_web'))

    return render_template('upload.html',now=datetime.now())

@app.route('/ui/delete-template-web', methods=['POST'])
@login_required
def delete_template_web():
    template = request.form.get('template')
    if template:
        path = os.path.join(app.config['UPLOAD_FOLDER'], template)
        try:
            os.remove(path)
            msg = f"UI Deleted Template: {template}"
            flash(msg) ; app.logger.info(msg)
        except Exception as e:
            msg = f"UI Error deleting Template {template}: {e}"
            flash(msg)
            app.logger.error(msg)
    else:
        flash("No template specified for deletion.")
    return redirect(url_for('list_templates_web'))



@app.route('/ui/placeholders', methods=['GET'])
@login_required
def ui_placeholders():
    template_name = request.args.get('template')
    if not template_name: return "Missing 'template' query parameter", 400

    template_path = os.path.join(app.config['UPLOAD_FOLDER'], template_name)
    if not os.path.exists(template_path): return f"Template '{template_name}' not found", 404

    slide_tag_map, slides_without_tags, all_tags = placeholder_mapper(template_path)

    return render_template(
    	'placeholders_table.html',
        template       = template_name,
        tags_by_slide  = slide_tag_map,
        slides_no_tags = slides_without_tags,
        all_tags       = all_tags,
        now            = datetime.now())

@app.route('/ui/metadata/<template>', methods=['GET'])
@login_required
def view_metadata(template):
    try:  metadata = get_metadata(template)
    except Exception as e:
        flash(f"Error reading metadata for {template}: {e}")
        return redirect(url_for('list_templates_web'))

    return render_template('view_metadata.html', template=template, metadata=metadata, now=datetime.now())


@app.route('/ui/fill-template-web', methods=['GET', 'POST'])
@login_required
def fill_template_web():
    template_name = request.args.get('template')
    if not template_name:  
        app.logger.error('UI Fill: Missing template=name.pptx')
        return "Missing 'template' query parameter", 400

    template_path = os.path.join(app.config['UPLOAD_FOLDER'], template_name)
    if not os.path.exists(template_path): 
    	app.logger.error(f'UI Fill: Could not find {template_name}')
    	return f"Template '{template_name}' not found", 404

    # RECEIVE RESPONSE
    if request.method == 'POST':
        filled_data = {key: request.form.get(key) for key in request.form}
        keep_slides = request.form.getlist('keep')

        try: 
            app.logger.info(f"UI Fill - Filling {template_name}")
            prs = fill_placeholders(template_path, filled_data)
        except Exception as e:
            app.logger.error(f"UI Fill: Failed to return placeholders: {e}")
            return "Error processing template", 500

        # FILTER SLIDEAS
        if keep_slides:
            try: filter_slides(prs, keep_slides, app.logger)
            except Exception as e: app.logger.warning(f"UI Fill: Slide filtering failed: {e}")

        # SAVE PROCESSED FILE
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        app.logger.info(f"UI Fill: Saving to Temp File: {temp_file.name}")
        prs.save(temp_file.name)
        temp_file.close()

        # SEND
        return send_file(temp_file.name, as_attachment  = True, download_name  = f"filled_{template_name}",
            mimetype  = 'application/vnd.openxmlformats-officedocument.presentationml.presentation')

    # BUILD THE WEBPAGE 
    try:  slide_tag_map, slides_without_tags, all_tags = placeholder_mapper(template_path)
    except Exception as e:
        app.logger.error(f"UI Fill: Failed to map placeholders: {e}")
        return "Error analyzing template", 500

    return render_template(
        'fill_template.html',
        template       = template_name,
        placeholders   = all_tags,
        tags_by_slide  = slide_tag_map,
        slides_no_tags = slides_without_tags,
        now            = datetime.now())


@app.route('/ui/api-key')
def show_api_key():
    api_key = app.config.get('API_KEY', 'Not configured')
    return render_template('show_api_key.html', api_key=api_key, now   = datetime.now())


#######################################################
#####  API ROUTES                                 #####
#######################################################

'''
https://p5-nvme.taild8d2e.ts.net/api/status?api_key=xxxxx
'''
@app.route('/api/status', methods=['GET'])
@require_api_key

def status():

    return jsonify({
        'service'             : 'ppt-template-service',
        'status'              : 'ok',
        'version'             : '0.2.1',
        'timestamp'           : datetime.utcnow().isoformat() + 'Z',
        'templates_available' : len(get_files()),
        'list_templates_url'  : url_for('list_templates', api_key=app.config['API_KEY'], _external=True),
        'ui_dashboard_url'    : url_for('ui_redirect', _external=True),
        'routes'              : get_routes(current_app),
    })


'''
curl -X POST http://localhost:5010/api/upload-template?api_key=xxxxx\
     -F "file=@demo.pptx"
'''

@app.route('/api/upload-template', methods=['POST'])
@require_api_key
def upload_template():
    file = request.files.get('file')
    if not file or not file.filename.endswith('.pptx'): 
        app.logger.error(f'API Upload: Invalid File Format for {file}')
        return jsonify({'error': 'Invalid file format'}), 400

    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)
    app.logger.info(f"API Upload: File saved successfully: {filepath}")


    return jsonify({'message': f'Template {filename} uploaded successfully'}), 200

'''
http://localhost:5010/api/download-template?template=demo.pptx&api_key=your-secret-key-here
'''

@app.route('/api/download-template', methods=['GET'])
@require_api_key
def download_template():
    template_name = request.args.get('template')
    template_path = os.path.join(app.config['UPLOAD_FOLDER'], template_name)

    if not os.path.exists(template_path): 
        app.logger.error(f"API Download: Could not find  {template_name}" )
        return jsonify({'error': 'Template not found'}), 404
    
    app.logger.info(f"API Downloading {template_name}" )
    return send_file(template_path, as_attachment=True)


'''
curl -X POST http://localhost:5010/api/remove-template?api_key=xxx \
     -H "Content-Type: application/json" \
     -d '{"template": "demo.pptx"}'

'''
@app.route('/api/remove-template', methods=['POST'])
@require_api_key
def remove_template():
    data = request.get_json()
    template_name = data.get('template')

    if not template_name: return jsonify({'error': "Missing 'template' field"}), 400

    template_path = os.path.join(app.config['UPLOAD_FOLDER'], template_name)

    if not os.path.exists(template_path): return jsonify({'error': f"Template '{template_name}' not found"}), 404

    try:
        os.remove(template_path)
        app.logger.warning(f"API Deleted template: {template_name}")
        return jsonify({'message': f"Template {template_name} deleted successfully"}), 200
    except Exception as e:
        app.logger.error(f"API Delete: Failed to delete template: {template_name} — {e}")
        return jsonify({'error': 'Failed to delete template'}), 500


'''
http://p5-nvme.taild8d2e.ts.net/api/list-templates?api_key=xxx 
'''
@app.route('/api/list-templates', methods=['GET'])
@require_api_key
def list_templates():
    templates = []
    for f in get_files():
        templates.append({
            'filename'         : f,
            'slide_count'      : slide_count(f),
            'placeholders_url': url_for('get_placeholders', template=f, api_key=app.config['API_KEY'], _external=True),
            'download_url': url_for('download_template', template=f, api_key=app.config['API_KEY'], _external=True),
        })

    return jsonify({'templates': templates,  'api_starting_url': url_for('status', api_key=app.config['API_KEY'], _external=True),})


'''
https://p5-nvme.taild8d2e.ts.net/api/placeholders?template=xxxxpptx?api_key=xxx 
'''

@app.route('/api/placeholders', methods=['GET'])
@require_api_key
def get_placeholders():

    template_name = request.args.get('template')
    if not template_name: return jsonify({'error': "Missing 'template' query parameter"}), 400

    template_path = os.path.join(app.config['UPLOAD_FOLDER'], template_name)
    if not os.path.exists(template_path): return jsonify({'error': 'Template not found'}), 404

    slide_tag_map, slides_without_tags, all_tags = placeholder_mapper(template_path)

    return jsonify({
        'template'        : template_name,
        'tags_by_slide'   : slide_tag_map,
        'slides_no_tags'  : slides_without_tags,
        'all_tags'        : all_tags,
        'list_templates'  : url_for('list_templates', api_key=app.config['API_KEY'], _external=True),


    })


'''
curl -X POST http://localhost:5000/api/fill-template?api_key=xxx \
  -H "Content-Type: application/json" \
  -d '{
    "template": "example.pptx",
    "replacements": {"KEY": "value",}' 
  --output filled_example.pptx

'''
@app.route('/api/fill-template', methods=['POST'])
@require_api_key
def fill_template():
    
    if request.is_json:
        data          = request.get_json()
        template_name = data.get('template')
        replacements  = data.get('replacements', {})
        slide_indices = data.get('keep')
        has_keep      = 'keep' in data
    else:
        template_name = request.form.get('template')
        replacements  = {key: request.form.get(key) for key in request.form if key != 'template'}
        slide_indices = request.form.getlist('keep')
        has_keep      = 'keep' in request.form

	# CHECK TEMPLATE KEY:NAME IS IN DATA
    if not template_name or not replacements:
        app.logger.error(f"API Fill: Missing Template or replacement keys in data")
        return jsonify({'error': 'Missing template name or replacements'}), 400

    # MAKE FULL PATH - check exists, return error if not
    template_path = os.path.join(app.config['UPLOAD_FOLDER'], template_name)
    if not os.path.exists(template_path): 
        app.logger.error(f"API Fill: Template not Found {template_name}")
        return jsonify({'error': 'Template not found'}), 404

    try: 
    	app.logger.info(f"API Fill: Processing {template_name}")
    	prs = fill_placeholders(template_path, replacements)
    	app.logger.info(f"API Fill: Replaced {len(replacements)} Tags in {template_name}")
    except Exception as e:
        app.logger.error(f"API Fill: Failed to fill placeholders: {e}")
        return jsonify({'error': 'Failed to process template'}), 500

    ### LIMIT UNWANTED SLIDES ###
    if has_keep and slide_indices:
        try: filter_slides(prs, slide_indices, app.logger)
        except Exception as e: app.logger.error(f"API Fill: Slide filtering failed: {e}")
    elif has_keep: app.logger.info("API Fill: 'keep' key provided but empty — keeping all slides")
    else: app.logger.info("API Fill: No 'keep' key provided — keeping all slides")

    ### SAVE FILE  ###
    filled_name = f"filled_{template_name}"
    filled_path = os.path.join(app.config['FILLED_FOLDER'], filled_name)
    
    try:
        prs.save(filled_path)
        app.logger.info(f"API Fill: Saved File {filled_name}")
    except Exception as e: app.logger.error(f"API Fill: Failed to Fill: {filled_path} — {e}")

	### DELETE FILE AFTER SENDING  ###
    if app.config.get('REMOVE_FILLED_FILES', False):
        @after_this_request
        def remove_file(response):
            try:
                os.remove(filled_path)
                app.logger.info(f"API Fill: Deleted cached file after sending: {filled_path}")
            except Exception as e:
                app.logger.warning(f"API Fill: Could not delete file {filled_path}: {e}")
            return response
    else:  app.logger.info(f"API Fill: Keeping Cached File: {filled_path}")

    app.logger.info(f"API Fill: Returning File: {filled_name}")

    return send_file(filled_path, as_attachment=True)


if __name__ == '__main__':
	app.run(host='0.0.0.0', port=5110, debug=True)




